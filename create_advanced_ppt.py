from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_advanced_ppt():
    prs = Presentation()
    
    # 16:9 Slide Size
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Professional Palette
    NAVY = RGBColor(26, 37, 47)      # #1a252f
    BLUE = RGBColor(52, 152, 219)    # #3498db
    BLUE_LIGHT = RGBColor(235, 245, 251)
    GRAY_LIGHT = RGBColor(248, 249, 250)
    WHITE = RGBColor(255, 255, 255)
    GOLD = RGBColor(241, 196, 15)
    TEXT_GRAY = RGBColor(80, 80, 80)

    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = GRAY_LIGHT

    def add_page_header(slide, title_text, slide_num):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
        rect.fill.solid()
        rect.fill.fore_color.rgb = NAVY
        rect.line.visible = False
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{slide_num}. {title_text}"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # --- Slide 1: Cover ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9), 0, Inches(4.333), prs.slide_height)
    acc.fill.solid()
    acc.fill.fore_color.rgb = NAVY
    acc.line.visible = False
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "BoGyeong Lim"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = "PORTFOLIO 2026"
    p2.font.size = Pt(32)
    p2.font.color.rgb = BLUE

    # --- Slide 2: ABOUT ME (With Photo) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_page_header(slide, "ABOUT ME", "01")
    
    # Photo Implementation
    photo_path = "portfolio-web/src/assets/profile.jpg"
    if os.path.exists(photo_path):
        # Add photo with slight border
        slide.shapes.add_picture(photo_path, Inches(1), Inches(1.5), Inches(3.8), Inches(5.1))
        # Add a decorative frame behind/around it
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.4), Inches(4), Inches(5.3))
        frame.fill.background()
        frame.line.color.rgb = BLUE
        frame.line.width = Pt(3)
    
    # Bio Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.5), Inches(7), Inches(5.1))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.visible = False
    
    txt_box = slide.shapes.add_textbox(Inches(5.8), Inches(1.8), Inches(6.4), Inches(4.5))
    tf = txt_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Focus: UX Flow & Maintainable Architecture"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = BLUE
    
    p = tf.add_paragraph()
    p.text = "단순 기능 구현을 넘어 상태 흐름의 명확성과 유지보수성을 최우선으로 생각합니다. 최근에는 오픈소스 생태계와 브라우저 엔진 기여에 깊은 관심을 가지고 연구 중입니다."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_GRAY
    p.space_before = Pt(20)
    
    for label, val in [("Phone", "010-3396-6856"), ("Email", "qhrudlim@gmail.com"), ("GitHub", "qhrudlim")]:
        p = tf.add_paragraph()
        p.text = f"{label}: {val}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.space_before = Pt(15)

    # --- Slide 3: TECH STACK (With Progress Bars) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_page_header(slide, "TECH STACK", "02")
    
    categories = [
        ("FRONTEND", [("React", 5), ("Vue", 4), ("HTML/CSS", 5)]),
        ("LANGUAGES", [("JS/TS", 5), ("Python", 4), ("Java", 2)]),
        ("BACKEND / DB", [("Django", 4), ("MySQL", 4)]),
        ("TOOLS", [("Git/GitHub", 5), ("Figma", 4), ("Jira/Postman", 4)])
    ]

    def add_skill_row(slide, x, y, name, level):
        # Skill Name
        st = slide.shapes.add_textbox(x, y, Inches(1.5), Inches(0.4))
        st.text_frame.paragraphs[0].text = name
        st.text_frame.paragraphs[0].font.size = Pt(14)
        st.text_frame.paragraphs[0].font.bold = True
        
        # 5 Gauge Boxes
        gauge_start_x = x + Inches(1.5)
        for i in range(5):
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, gauge_start_x + Inches(i * 0.35), y + Inches(0.05), Inches(0.3), Inches(0.2))
            box.line.visible = False
            if i < level:
                box.fill.solid()
                box.fill.fore_color.rgb = BLUE
            else:
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(220, 220, 220)

    for i, (cat_name, skills) in enumerate(categories):
        x = Inches(0.5 + (i * 3.2))
        y_top = Inches(1.5)
        
        # Category Box
        cbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_top, Inches(3), Inches(5))
        cbox.fill.solid()
        cbox.fill.fore_color.rgb = WHITE
        cbox.line.color.rgb = BLUE_LIGHT
        
        # Category Title
        ct = slide.shapes.add_textbox(x, y_top + Inches(0.2), Inches(3), Inches(0.5))
        ct.text_frame.paragraphs[0].text = cat_name
        ct.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        ct.text_frame.paragraphs[0].font.bold = True
        ct.text_frame.paragraphs[0].font.size = Pt(18)
        ct.text_frame.paragraphs[0].font.color.rgb = NAVY
        
        # Add Skill Rows
        for idx, (sname, slevel) in enumerate(skills):
            add_skill_row(slide, x + Inches(0.1), y_top + Inches(1 + (idx * 0.8)), sname, slevel)

    # --- Slide 4: MILESTONES (Clean Timeline) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_page_header(slide, "MILESTONES", "03")
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(1.5), Inches(0.05), Inches(5.5))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.visible = False
    
    milestones = [
        ("2026.06", "SSAFY 14기 수료", "구미캠퍼스 | 공통 프로젝트 반 우승 수상"),
        ("2025.08", "충남대 석사 졸업", "과학수사학 전공"),
        ("2025.06", "웹디자인 교육 수료", "SBS 아카데미 컴퓨터아트학원"),
        ("2022.02", "계명대 학사 졸업", "화학 전공"),
        ("2016.02", "한국사 2급 취득", "국사편찬위원회")
    ]
    
    for i, (date, title, desc) in enumerate(milestones):
        y = Inches(1.8 + (i * 1.1))
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.375), y, Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.color.rgb = WHITE
        
        dt = slide.shapes.add_textbox(Inches(0.5), y, Inches(1.8), Inches(0.4))
        dt.text_frame.paragraphs[0].text = date
        dt.text_frame.paragraphs[0].font.bold = True
        dt.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        ct = slide.shapes.add_textbox(Inches(3), y - Inches(0.1), Inches(9.5), Inches(0.8))
        tf = ct.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = NAVY
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY

    # --- Slide 5: FEATURED PROJECT - DocQ ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_page_header(slide, "FEATURED PROJECT", "04")
    
    title_area = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(1))
    tf = title_area.text_frame
    p = tf.paragraphs[0]
    p.text = "DocQ: PDF-based Multiplayer 3D Quiz Game"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    p = tf.add_paragraph()
    p.text = "🏆 SSAFY 공통 프로젝트 반 우승 수상 | Frontend Lead"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    # Content Columns
    left_txt = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(6), Inches(4))
    tf = left_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Challenges & Solutions"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = NAVY
    
    for item in [
        "• 2.5D 한계 극복 -> 3D 통합 구조 전환 주도",
        "• GLTF 최적화를 통한 FBX 에셋 로딩 시간 단축",
        "• 상태 전이 기반 자동 진행 게임 로직 설계",
        "• 멀티플레이 환경 단계 동기화 구조 정립"
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_before = Pt(12)

    # Right side visual box
    vis_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2.5), Inches(5.8), Inches(4))
    vis_box.fill.solid()
    vis_box.fill.fore_color.rgb = WHITE
    vis_box.line.color.rgb = RGBColor(200, 200, 200)
    vis_box.line.dash_style = 2
    vt = slide.shapes.add_textbox(Inches(7.5), Inches(4), Inches(5), Inches(1))
    vt.text_frame.paragraphs[0].text = "[ Project Screenshot Area ]"
    vt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Save
    prs.save("portfolio-11.pptx")
    print("Success: portfolio-11.pptx created with photo and skill gauges.")

if __name__ == "__main__":
    create_advanced_ppt()
