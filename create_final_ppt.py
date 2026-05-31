from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_ppt():
    prs = Presentation()
    
    # 16:9 Slide Size
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    PRIMARY_COLOR = RGBColor(44, 62, 80)    # #2c3e50
    ACCENT_COLOR = RGBColor(52, 152, 219)  # #3498db
    BG_COLOR = RGBColor(248, 249, 250)      # #f8f9fa

    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_title_header(slide, title_text, slide_num):
        # Top Bar
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY_COLOR
        shape.line.visible = False
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(10), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{slide_num}. {title_text}"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

    # --- Slide 1: Cover ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    
    # Center Rectangle
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(9.33), Inches(3.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = PRIMARY_COLOR
    rect.line.visible = False
    
    # Text
    title_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.5), Inches(8.33), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PORTFOLIO"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    subtitle_box = slide.shapes.add_textbox(Inches(2.5), Inches(4), Inches(8.33), Inches(0.8))
    p = subtitle_box.text_frame.paragraphs[0]
    p.text = "qhrudlim | Frontend Developer 임보경"
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_COLOR

    # --- Slide 2: ABOUT ME ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title_header(slide, "ABOUT ME", "01")
    
    # Philosophy Box
    philosophy = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(1))
    p = philosophy.text_frame.paragraphs[0]
    p.text = '"사용자 경험 흐름을 설계하고 단순하고 유지보수하기 쉬운 구조를 유지하는 데 중점을 둡니다."'
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = PRIMARY_COLOR

    # Bio Text
    bio = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(6), Inches(3))
    tf = bio.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "단순한 기능 구현을 넘어 상태 흐름의 명확성과 유지보수성을 최우선으로 생각하는 프론트엔드 개발자입니다. 문제를 구조화하여 정리하는 것을 즐기며, UX 고도화 과정에서 보람을 느낍니다. 최근에는 오픈소스 생태계와 브라우저 엔진 기여에 깊은 관심을 가지고 공부하고 있습니다."
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(85, 85, 85)

    # Contact & Personal Info
    contact = slide.shapes.add_textbox(Inches(7.5), Inches(2.5), Inches(5), Inches(2))
    tf = contact.text_frame
    for line in ["Name: 임보경 (BoGyeong Lim)", "Phone: 010-3396-6856", "Email: qhrudlim@gmail.com", "GitHub: qhrudlim"]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR

    # --- Slide 3: TECH STACK ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title_header(slide, "TECH STACK", "02")
    
    # Grid of Boxes for Categories
    stacks = [
        ("Frontend", "React, Vue, HTML, Vanilla CSS, Tailwind CSS"),
        ("Languages", "JavaScript, TypeScript, Python, Java (Learning)"),
        ("Backend & DB", "Django, MySQL"),
        ("Tools", "Git, GitHub, GitLab, Jira, Figma, Mattermost, Postman")
    ]
    
    x_pos = [Inches(0.5), Inches(6.8)]
    y_pos = [Inches(1.2), Inches(4.2)]
    
    for i, (cat, items) in enumerate(stacks):
        x = x_pos[i % 2]
        y = y_pos[i // 2]
        
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = ACCENT_COLOR
        
        # Text
        txt = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(5.6), Inches(2))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = cat
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        p = tf.add_paragraph()
        p.text = items
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(68, 68, 68)

    # --- Slide 4: TIMELINE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title_header(slide, "TIMELINE & HISTORY", "03")
    
    timeline_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(5.5))
    tf = timeline_box.text_frame
    tf.word_wrap = True
    
    history = [
        "2026.06 | SSAFY 14기 수료 (구미캠퍼스) - 공통 프로젝트 반 우승",
        "2025.08 | 충남대학교 과학수사학과 석사 졸업",
        "2025.07 | 삼성 청년 SW 아카데미(SSAFY) 입과",
        "2025.06 | SBS 아카데미 웹디자인 과정 수료",
        "2022.11 | 과학 강사 (와이즈만영재교육 대전도안센터)",
        "2022.02 | 계명대학교 화학전공 학사 졸업",
        "2016.02 | 한국사능력검정시험 2급 취득"
    ]
    
    for item in history:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(68, 68, 68)
        p.space_after = Pt(12)

    # --- Slide 5: PROJECT - DocQ ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title_header(slide, "PROJECT: DocQ", "04")
    
    # Award Badge
    badge = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(10.5), Inches(1.2), Inches(2), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(241, 196, 15)
    badge.line.visible = False
    bt = slide.shapes.add_textbox(Inches(10.5), Inches(1.2), Inches(2), Inches(0.6))
    bt.text_frame.paragraphs[0].text = "반 우승 수상"
    bt.text_frame.paragraphs[0].font.bold = True
    bt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    lines = [
        "제목: DocQ - PDF 기반 멀티플레이 3D 퀴즈 게임",
        "역할: 프론트엔드 (3D 게임 화면 및 상태 기반 흐름 설계)",
        "핵심 기술: React, Three.js, TypeScript, Blender",
        "",
        "주요 성과 및 해결 과정:",
        "1. 2.5D 한계 극복을 위한 3D 통합 구조 전환 주도",
        "2. Blender를 활용한 GLTF/GLB 포맷 변환으로 로딩 시간 최적화",
        "3. 상태 전이 기반 자동 진행 게임 로직 설계 및 구현"
    ]
    
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        if "성과" in line: p.font.bold = True

    # --- Slide 6: PROJECTS - BioTwin & JIPCHAK ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title_header(slide, "OTHER PROJECTS", "05")
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5))
    tf = content.text_frame
    
    projects = [
        ("BioTwin (특화 프로젝트)", "• 생체 데이터 시각화 플랫폼", "• Vue.js, TypeScript 기반 데이터 시각화 차트 구현"),
        ("", "", ""),
        ("JIPCHAK (자율 프로젝트)", "• IoT 기반 스마트 솔루션", "• React, Tailwind CSS, TypeScript 활용 메인 페이지 구현", "• 하드웨어 제작 프로세스 및 오픈소스 기여")
    ]
    
    for title, desc, detail, *extra in projects:
        if title:
            p = tf.add_paragraph()
            p.text = title
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = ACCENT_COLOR
        if desc:
            p = tf.add_paragraph()
            p.text = desc
            p.font.size = Pt(16)
        if detail:
            p = tf.add_paragraph()
            p.text = detail
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(102, 102, 102)

    prs.save("final_portfolio_qhrudlim.pptx")
    print("PPT created successfully: final_portfolio_qhrudlim.pptx")

if __name__ == "__main__":
    create_ppt()
