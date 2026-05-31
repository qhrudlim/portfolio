from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def set_text_preserve_style(shape, new_text_list):
    """Replaces text in a shape paragraph by paragraph to keep existing run styles if possible."""
    if not hasattr(shape, "text_frame"):
        return
    
    tf = shape.text_frame
    # If new_text_list is a single string, wrap it
    if isinstance(new_text_list, str):
        new_text_list = [new_text_list]
    
    # Try to reuse existing paragraphs and runs to keep fonts
    for i, text in enumerate(new_text_list):
        if i < len(tf.paragraphs):
            p = tf.paragraphs[i]
            if p.runs:
                p.runs[0].text = text
                # Clear other runs in this paragraph
                for r_idx in range(1, len(p.runs)):
                    p.runs[r_idx].text = ""
            else:
                p.text = text
        else:
            p = tf.add_paragraph()
            p.text = text

def create_final_portfolio():
    prs = Presentation("portfolio-5.pptx")
    
    # Slide 1: Title (Keep)
    
    # Slide 2: PROFILE
    # Shape 13 is the main intro text box
    set_text_preserve_style(prs.slides[1].shapes[12], [
        "사용자 경험의 흐름을 논리적으로 설계하고,",
        "지속 가능한 코드 구조를 지향하는 프론트엔드 개발자 임보경입니다.",
        "단순 구현보다 상태 흐름의 명확성과 유지보수 가능한 구조를 중요하게 생각합니다."
    ])

    # Slide 3: TECH STACK
    # Shape 15 is Tools
    set_text_preserve_style(prs.slides[2].shapes[14], [
        "• Git",
        "• Jira",
        "• Figma",
        "• Mattermost"
    ])

    # Slide 4: 03. PROJECT: DocQ (Keep)
    # Slide 5: 04. TECHNICAL CHALLENGES (Keep)

    # Slide 6: 05. OTHER PROJECTS (New)
    # We use Slide 4 as a visual template but build it shape by shape to be sure.
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)
    
    # Background (Rectangle 1)
    bg = new_slide.shapes.add_shape(1, 0, 0, 9144000, 6858000)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF1, 0xF2, 0xF6)
    bg.line.visible = False
    
    # Top Bar (Rectangle 2)
    top_bar = new_slide.shapes.add_shape(1, 0, 0, 9144000, 914400)
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    top_bar.line.visible = False
    
    # Title (TextBox 3)
    title_box = new_slide.shapes.add_textbox(457200, 182880, 8229600, 548640)
    title_box.text_frame.text = "05. OTHER PROJECTS"
    # Try to match font from Slide 4 title
    source_title = prs.slides[3].shapes[2]
    if source_title.text_frame.paragraphs[0].runs:
        source_run = source_title.text_frame.paragraphs[0].runs[0]
        title_run = title_box.text_frame.paragraphs[0].runs[0]
        title_run.font.name = source_run.font.name
        title_run.font.size = source_run.font.size
        title_run.font.bold = source_run.font.bold
        if source_run.font.color and hasattr(source_run.font.color, 'rgb'):
            title_run.font.color.rgb = source_run.font.color.rgb

    # Content Box (Rectangle 4)
    content_box = new_slide.shapes.add_shape(1, 457200, 1371600, 8229600, 5029200)
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    content_box.line.visible = False
    
    # Content Text (TextBox 5 & 6 combined)
    text_box = new_slide.shapes.add_textbox(731520, 1645920, 7680960, 4500000)
    tf = text_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "BioTwin & JI[1]PCHAK"
    p.font.bold = True
    p.font.size = Pt(20)
    
    tf.add_paragraph() # Spacer
    
    p = tf.add_paragraph()
    p.text = "BioTwin (SSAFY 특화 프로젝트)"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• 디지털 트윈 기술 기반 도메인 프로젝트 수행"
    
    tf.add_paragraph() # Spacer
    
    p = tf.add_paragraph()
    p.text = "JI[1]PCHAK (SSAFY 자율 프로젝트)"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• 메인 페이지 및 GUIDE 구현"
    p = tf.add_paragraph()
    p.text = "• 하드웨어 제작 참여 및 오픈소스 프로젝트"

    # Save
    prs.save("portfolio-6.pptx")
    print("Successfully created portfolio-6.pptx with matched style and correct numbering.")

if __name__ == "__main__":
    create_final_portfolio()
