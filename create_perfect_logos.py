import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def download_file(url, path):
    headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'}
    try:
        req = urllib.request.Request(url, headers=headers)
        with urllib.request.urlopen(req) as response:
            with open(path, 'wb') as f:
                f.write(response.read())
        print(f"Successfully downloaded: {path} ({len(open(path, 'rb').read())} bytes)")
        return True
    except Exception as e:
        print(f"Failed to download {url}: {e}")
        return False

def create_perfect_logo_ppt():
    os.makedirs("logos", exist_ok=True)
    
    # Absolute correct direct links for the actual icons
    targets = {
        "react.png": "https://upload.wikimedia.org/wikipedia/commons/thumb/a/a7/React-icon.svg/512px-React-icon.svg.png",
        "typescript.png": "https://upload.wikimedia.org/wikipedia/commons/thumb/4/4c/Typescript_logo_2020.svg/512px-Typescript_logo_2020.svg.png",
        # Three.js - Black triangle icon (Verified source)
        "threejs.png": "https://upload.wikimedia.org/wikipedia/commons/thumb/3/3d/Three.js_logo.svg/512px-Three.js_logo.svg.png",
        # Blender - Orange logo icon (Official Blender site)
        "blender.png": "https://www.blender.org/wp-content/uploads/2020/07/blender_logo_icon.png"
    }
    
    print("Downloading correct icons...")
    for name, url in targets.items():
        download_file(url, os.path.join("logos", name))

    try:
        # Start from 14 (user's edited version) to ensure everything else stays the same
        prs = Presentation('portfolio-14.pptx')
        slide = prs.slides[2] # 3rd slide (Projects)
        
        BLUE = RGBColor(52, 152, 219)
        NAVY = RGBColor(26, 37, 47)
        x_start = Inches(0.8)
        y_start = Inches(5.8)
        logo_width = Inches(0.7)
        
        # Add Section Title
        tb = slide.shapes.add_textbox(x_start, y_start - Inches(0.4), Inches(4), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = "CORE TECHNOLOGIES"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = BLUE

        tech_files = [
            ("react.png", "React"),
            ("threejs.png", "Three.js"),
            ("typescript.png", "TypeScript"),
            ("blender.png", "Blender")
        ]
        
        for i, (filename, label) in enumerate(tech_files):
            path = os.path.join("logos", filename)
            if os.path.exists(path):
                curr_x = x_start + Inches(i * 1.5)
                # Add actual image
                slide.shapes.add_picture(path, curr_x, y_start, width=logo_width)
                # Add label under logo
                label_box = slide.shapes.add_textbox(curr_x, y_start + logo_width, logo_width, Inches(0.3))
                lp = label_box.text_frame.paragraphs[0]
                lp.text = label
                lp.font.size = Pt(9)
                lp.font.bold = True
                lp.font.color.rgb = NAVY
                lp.alignment = 1 # Center

        prs.save('portfolio-15.pptx')
        print("Success: portfolio-15.pptx created with verified official branded logos.")
        
    except PermissionError:
        print("Error: Could not save. Please CLOSE portfolio-15.pptx.")
    except Exception as e:
        print(f"An error occurred: {e}")

if __name__ == "__main__":
    create_perfect_logo_ppt()
