import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def download_file(url, path):
    headers = {'User-Agent': 'Mozilla/5.0'}
    try:
        req = urllib.request.Request(url, headers=headers)
        with urllib.request.urlopen(req) as response:
            with open(path, 'wb') as f:
                f.write(response.read())
        print(f"Successfully downloaded: {path}")
        return True
    except Exception as e:
        print(f"Failed {url}: {e}")
        return False

def create_perfect_logo_ppt_v3():
    os.makedirs("logos", exist_ok=True)
    
    # Final verified direct links
    targets = {
        "react.png": "https://raw.githubusercontent.com/github/explore/80688e429a7d4ef2fca1e82350fe8e3517d3494d/topics/react/react.png",
        "typescript.png": "https://raw.githubusercontent.com/github/explore/80688e429a7d4ef2fca1e82350fe8e3517d3494d/topics/typescript/typescript.png",
        "threejs.png": "https://raw.githubusercontent.com/github/explore/80688e429a7d4ef2fca1e82350fe8e3517d3494d/topics/threejs/threejs.png",
        "blender.png": "https://img.icons8.com/color/512/blender-3d.png"
    }
    
    for name, url in targets.items():
        download_file(url, os.path.join("logos", name))

    try:
        prs = Presentation('portfolio-14.pptx')
        slide = prs.slides[2] 
        BLUE = RGBColor(52, 152, 219)
        NAVY = RGBColor(26, 37, 47)
        x_start = Inches(0.8)
        y_start = Inches(5.8)
        logo_width = Inches(0.7)
        
        tb = slide.shapes.add_textbox(x_start, y_start - Inches(0.4), Inches(4), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = "CORE TECHNOLOGIES"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = BLUE

        tech_files = [
            ("react.png", "React"),
            ("threejs.png", "Three.js"),
            ("typescript.png", "TypeScript"),
            ("blender.png", "Blender")
        ]
        
        for i, (filename, label) in enumerate(tech_files):
            path = os.path.join("logos", filename)
            if os.path.exists(path):
                curr_x = x_start + Inches(i * 1.4)
                slide.shapes.add_picture(path, curr_x, y_start, width=logo_width)
                label_box = slide.shapes.add_textbox(curr_x, y_start + logo_width, logo_width, Inches(0.3))
                lp = label_box.text_frame.paragraphs[0]
                lp.text = label
                lp.font.size = Pt(9)
                lp.font.bold = True
                lp.font.color.rgb = NAVY
                lp.alignment = 1

        prs.save('portfolio-15.pptx')
        print("Success: portfolio-15.pptx updated with correct official icons.")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    create_perfect_logo_ppt_v3()
