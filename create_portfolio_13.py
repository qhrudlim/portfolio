from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_compact_ppt():
    prs = Presentation()
    
    # 16:9 Slide Size
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Palette from Version 10
    NAVY = RGBColor(26, 37, 47)      
    BLUE = RGBColor(52, 152, 219)    
    GRAY_LIGHT = RGBColor(241, 242, 246) 
    WHITE = RGBColor(255, 255, 255)
    GOLD = RGBColor(241, 196, 15)
    TEXT_GRAY = RGBColor(80, 80, 80)

    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = GRAY_LIGHT

    def add_page_header(slide, title_text, slide_num):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        rect.fill.solid()
        rect.fill.fore_color.rgb = NAVY
        rect.line.visible = False
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), prs.slide_width, Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = BLUE
        line.line.visible = False

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{slide_num}. {title_text}"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT

    # --- Slide 1: Cover ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.333), 0, Inches(5), prs.slide_height)
    acc.fill.solid()
    acc.fill.fore_color.rgb = NAVY
    acc.line.visible = False
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(7), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "BoGyeong Lim"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = "PORTFOLIO 2026"
    p2.font.size = Pt(28)
    p2.font.color.rgb = BLUE

    # --- Slide 2: ABOUT & SKILLS (Merged Page) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_page_header(slide, "ABOUT & SKILLS", "01")
    
    # Left Column: Profile (0.5 to 5.0)
    # Smaller Photo
    photo_path = "portfolio-web/src/assets/profile.jpg"
    photo_w = Inches(2.2)
    photo_h = Inches(2.93) # 3:4 ratio
    if os.path.exists(photo_path):
        slide.shapes.add_picture(photo_path, Inches(0.5), Inches(1.5), photo_w, photo_h)
    
    # Condensed Bio
    bio_txt = slide.shapes.add_textbox(Inches(2.9), Inches(1.5), Inches(2.4), Inches(2.93))
    tf = bio_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "BoGyeong Lim"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = NAVY
    
    p = tf.add_paragraph()
    p.text = "단순 기능 구현을 넘어 상태 흐름의 명확성과 유지보수성을 최우선으로 생각하는 프론트엔드 개발자입니다."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_GRAY
    p.space_before = Pt(5)

    # Single Contact Card
    contact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.6), Inches(4.8), Inches(2.3))
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = WHITE
    contact_box.line.color.rgb = BLUE
    
    ct = slide.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(4.4), Inches(2))
    tf = ct.text_frame
    p = tf.paragraphs[0]
    p.text = "CONTACT"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE
    
    for label, val in [("Phone", "010-3396-6856"), ("Email", "qhrudlim@gmail.com"), ("GitHub", "qhrudlim")]:
        p = tf.add_paragraph()
        p.text = f"• {label}: {val}"
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY
        p.space_before = Pt(5)

    # Right Column: Tech Stack Grid (5.8 to 12.8)
    stacks = [
        ("FRONTEND", [("React", 5), ("Vue", 4), ("HTML/CSS", 5)], BLUE),
        ("LANGUAGES", [("JS/TS", 5), ("Python", 4), ("Java", 2)], NAVY),
        ("BACKEND / DB", [("Django", 4), ("MySQL", 4)], BLUE),
        ("TOOLS", [("Git/GitHub", 5), ("Figma", 4), ("Jira", 4)], NAVY)
    ]
    
    def add_gauge(slide, x, y, level):
        for i in range(5):
            gx = x + Inches(i * 0.18)
            gs = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, gx, y + Inches(0.04), Inches(0.14), Inches(0.14))
            gs.line.visible = False
            if i < level:
                gs.fill.solid()
                gs.fill.fore_color.rgb = WHITE
            else:
                gs.fill.solid()
                gs.fill.fore_color.rgb = RGBColor(180, 180, 180)

    # 2x2 Grid for right side
    for i, (cat, items, color) in enumerate(stacks):
        grid_x = Inches(5.6 + (i % 2) * 3.7)
        grid_y = Inches(1.5 + (i // 2) * 2.8)
        
        # Skill Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, grid_x, grid_y, Inches(3.5), Inches(2.6))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.visible = False
        
        # Title
        tb = slide.shapes.add_textbox(grid_x, grid_y + Inches(0.1), Inches(3.5), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = cat
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Skills
        for idx, (sname, slevel) in enumerate(items):
            sy = grid_y + Inches(0.7 + (idx * 0.55))
            nt = slide.shapes.add_textbox(grid_x + Inches(0.15), sy, Inches(2), Inches(0.3))
            p = nt.text_frame.paragraphs[0]
            p.text = sname
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            add_gauge(slide, grid_x + Inches(2.2), sy + Inches(0.02), slevel)

    # --- Slide 3: MILESTONES (Version 10 Style) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_page_header(slide, "MILESTONES", "02")
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(1.5), Inches(0.05), Inches(5.5))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.visible = False
    
    milestones = [
        ("2026.06", "SSAFY 14기 수료", "구미캠퍼스 | 공통 프로젝트 반 우승"),
        ("2025.08", "충남대 석사 졸업", "과학수사학 전공"),
        ("2025.06", "웹디자인 교육 수료", "SBS 아카데미 컴퓨터아트학원"),
        ("2022.02", "계명대 학사 졸업", "화학 전공"),
        ("2016.02", "한국사 2급 취득", "국사편찬위원회")
    ]
    for i, (date, title, desc) in enumerate(milestones):
        y = Inches(1.8 + (i * 1.1))
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.85), y, Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.color.rgb = WHITE
        dt = slide.shapes.add_textbox(Inches(0.5), y, Inches(1.2), Inches(0.4))
        dt.text_frame.paragraphs[0].text = date
        dt.text_frame.paragraphs[0].font.bold = True
        dt.text_frame.paragraphs[0].font.color.rgb = BLUE
        cb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), y - Inches(0.1), Inches(10), Inches(0.8))
        cb.fill.solid()
        cb.fill.fore_color.rgb = WHITE
        cb.line.visible = False
        ct = slide.shapes.add_textbox(Inches(2.7), y, Inches(9.5), Inches(0.6))
        p = ct.text_frame.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        p = ct.text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY

    # --- Slide 4: FEATURED PROJECT ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_page_header(slide, "FEATURED PROJECT: DocQ", "03")
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.width = Pt(2)
    card.line.color.rgb = BLUE
    sum_txt = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11), Inches(1))
    p = sum_txt.text_frame.paragraphs[0]
    p.text = "DocQ: PDF-based Multiplayer 3D Quiz Game"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = NAVY
    p = sum_txt.text_frame.add_paragraph()
    p.text = "🏆 SSAFY 공통 프로젝트 반 우승 수상 | 2026.01.06 ~ 2026.02.13 (6주)"
    p.font.bold = True
    p.font.color.rgb = GOLD
    cl_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.2), Inches(6), Inches(3.5))
    cl_box.fill.solid()
    cl_box.fill.fore_color.rgb = WHITE
    cl_box.line.visible = False
    cl_txt = slide.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(5.6), Inches(3))
    p = cl_txt.text_frame.paragraphs[0]
    p.text = "Technical Challenges"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = BLUE
    for issue in ["• 2.5D 렌더링 한계 -> 3D 통합 구조 전환", "• FBX 에셋 로딩 문제 -> GLTF 최적화", "• 상태 기반 게임 흐름 자동화 설계"]:
        p = cl_txt.text_frame.add_paragraph()
        p.text = issue
        p.font.size = Pt(14)
        p.space_before = Pt(10)
    im_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.2), Inches(6), Inches(3.5))
    im_box.fill.solid()
    im_box.fill.fore_color.rgb = WHITE
    im_box.line.color.rgb = RGBColor(200, 200, 200)
    im_box.line.dash_style = 2
    im_txt = slide.shapes.add_textbox(Inches(7.2), Inches(4.5), Inches(5), Inches(1))
    im_txt.text_frame.paragraphs[0].text = "Project Screenshot Area"
    im_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    im_txt.text_frame.paragraphs[0].font.color.rgb = RGBColor(180, 180, 180)

    # --- Slide 5: ADDITIONAL PROJECTS ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_page_header(slide, "ADDITIONAL PROJECTS", "04")
    projs = [
        ("BioTwin (특화)", "생체 데이터 시각화 플랫폼", "Vue.js / TypeScript", "실시간 모니터링 대시보드 및 차트 구현"),
        ("JIPCHAK (자율)", "IoT 스마트 솔루션", "React / Tailwind / IoT", "메인 페이지 UI/UX 및 하드웨어 연동")
    ]
    for i, (name, sub, tech, detail) in enumerate(projs):
        x = Inches(0.5 + (i * 6.5))
        y = Inches(1.8)
        pc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.8), Inches(4.5))
        pc.fill.solid()
        pc.fill.fore_color.rgb = WHITE
        pc.line.color.rgb = BLUE
        pt = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), Inches(5.2), Inches(4))
        tf = pt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = NAVY
        p = tf.add_paragraph()
        p.text = sub
        p.font.size = Pt(16)
        p.font.color.rgb = BLUE
        p = tf.add_paragraph()
        p.text = f"Stack: {tech}"
        p.font.bold = True
        p.font.size = Pt(14)
        p.space_before = Pt(15)
        p = tf.add_paragraph()
        p.text = detail
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        p.space_before = Pt(10)

    # Save
    prs.save("portfolio-13.pptx")
    print("Success: portfolio-13.pptx created with merged About & Skills and optimized layout.")

if __name__ == "__main__":
    create_compact_ppt()
