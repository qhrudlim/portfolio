from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_final_slogan_ppt():
    prs = Presentation()
    
    # 16:9 Slide Size
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Palette from Version 10
    NAVY = RGBColor(26, 37, 47)      
    BLUE = RGBColor(52, 152, 219)    
    GRAY_LIGHT = RGBColor(241, 242, 246) 
    WHITE = RGBColor(255, 255, 255)
    GOLD = RGBColor(241, 196, 15)
    TEXT_GRAY = RGBColor(80, 80, 80)

    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = GRAY_LIGHT

    def add_page_header(slide, title_text, slide_num):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        rect.fill.solid()
        rect.fill.fore_color.rgb = NAVY
        rect.line.visible = False
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), prs.slide_width, Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = BLUE
        line.line.visible = False

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{slide_num}. {title_text}"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT

    # --- Slide 1: Cover (With Slogan) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.333), 0, Inches(5), prs.slide_height)
    acc.fill.solid()
    acc.fill.fore_color.rgb = NAVY
    acc.line.visible = False
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(7), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "BoGyeong Lim"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    p2 = tf.add_paragraph()
    p2.text = "논리적 구조로 단단하게, 직관적 흐름으로 유연하게"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = BLUE
    
    p3 = tf.add_paragraph()
    p3.text = "매일 더 나은 코드를 탐구하는 프론트엔드 설계자"
    p3.font.size = Pt(18)
    p3.font.color.rgb = TEXT_GRAY

    # --- Slide 2: ABOUT & SKILLS (Merged Page + Slogan & 3 Pillars) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_page_header(slide, "ABOUT & SKILLS", "01")
    
    # Smaller Photo
    photo_path = "portfolio-web/src/assets/profile.jpg"
    if os.path.exists(photo_path):
        slide.shapes.add_picture(photo_path, Inches(0.5), Inches(1.5), Inches(2.2), Inches(2.93))
    
    # Bio & Slogan Area
    bio_txt = slide.shapes.add_textbox(Inches(2.9), Inches(1.5), Inches(2.5), Inches(2.93))
    tf = bio_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '"논리적 구조로 단단하게, \n직관적 흐름으로 유연하게"'
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = BLUE
    
    # 3 Pillars
    pillars = [
        ("Structure", "논리적인 구조 설계를 통해 유지보수가 용이한 코드를 지향합니다."),
        ("UX Flow", "사용자의 동선을 고려하여 직관적이고 매끄러운 경험을 제공합니다."),
        ("Growth", "새로운 기술에 대한 깊은 탐구와 리팩토링을 통해 성장을 추구합니다.")
    ]
    
    for title, desc in pillars:
        p = tf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = NAVY
        p.space_before = Pt(8)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_GRAY

    # Contact Card
    contact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.8), Inches(4.8), Inches(2.1))
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = WHITE
    contact_box.line.color.rgb = BLUE
    ct = slide.shapes.add_textbox(Inches(0.7), Inches(4.9), Inches(4.4), Inches(1.8))
    p = ct.text_frame.paragraphs[0]
    p.text = "CONTACT"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = BLUE
    for label, val in [("Phone", "010-3396-6856"), ("Email", "qhrudlim@gmail.com"), ("GitHub", "qhrudlim")]:
        p = ct.text_frame.add_paragraph()
        p.text = f"• {label}: {val}"
        p.font.size = Pt(11)
        p.font.color.rgb = NAVY

    # Right Column: Tech Stack Grid (No Changes)
    stacks = [
        ("FRONTEND", [("React", 5), ("Vue", 4), ("HTML/CSS", 5)], BLUE),
        ("LANGUAGES", [("JS/TS", 5), ("Python", 4), ("Java", 2)], NAVY),
        ("BACKEND / DB", [("Django", 4), ("MySQL", 4)], BLUE),
        ("TOOLS", [("Git/GitHub", 5), ("Figma", 4), ("Jira", 4)], NAVY)
    ]
    def add_gauge(slide, x, y, level):
        for i in range(5):
            gx = x + Inches(i * 0.18)
            gs = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, gx, y + Inches(0.04), Inches(0.14), Inches(0.14))
            gs.line.visible = False
            if i < level:
                gs.fill.solid()
                gs.fill.fore_color.rgb = WHITE
            else:
                gs.fill.solid()
                gs.fill.fore_color.rgb = RGBColor(180, 180, 180)

    for i, (cat, items, color) in enumerate(stacks):
        grid_x = Inches(5.6 + (i % 2) * 3.7)
        grid_y = Inches(1.5 + (i // 2) * 2.8)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, grid_x, grid_y, Inches(3.5), Inches(2.6))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.visible = False
        tb = slide.shapes.add_textbox(grid_x, grid_y + Inches(0.1), Inches(3.5), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = cat
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        for idx, (sname, slevel) in enumerate(items):
            sy = grid_y + Inches(0.7 + (idx * 0.55))
            nt = slide.shapes.add_textbox(grid_x + Inches(0.15), sy, Inches(2), Inches(0.3))
            p = nt.text_frame.paragraphs[0]
            p.text = sname
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            add_gauge(slide, grid_x + Inches(2.2), sy + Inches(0.02), slevel)

    # --- Slide 3: MILESTONES (No Changes) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_page_header(slide, "MILESTONES", "02")
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(1.5), Inches(0.05), Inches(5.5))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.visible = False
    milestones = [("2026.06", "SSAFY 14기 수료", "구미캠퍼스 | 공통 프로젝트 반 우승"), ("2025.08", "충남대 석사 졸업", "과학수사학 전공"), ("2025.06", "웹디자인 교육 수료", "SBS 아카데미"), ("2022.02", "계명대 학사 졸업", "화학 전공"), ("2016.02", "한국사 2급 취득", "국사편찬위원회")]
    for i, (date, title, desc) in enumerate(milestones):
        y = Inches(1.8 + (i * 1.1))
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.85), y, Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.color.rgb = WHITE
        dt = slide.shapes.add_textbox(Inches(0.5), y, Inches(1.2), Inches(0.4))
        dt.text_frame.paragraphs[0].text = date
        dt.text_frame.paragraphs[0].font.bold = True
        dt.text_frame.paragraphs[0].font.color.rgb = BLUE
        cb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), y - Inches(0.1), Inches(10), Inches(0.8))
        cb.fill.solid()
        cb.fill.fore_color.rgb = WHITE
        cb.line.visible = False
        ct = slide.shapes.add_textbox(Inches(2.7), y, Inches(9.5), Inches(0.6))
        p = ct.text_frame.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        p = ct.text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY

    # --- Slide 4: FEATURED PROJECT (No Changes) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_page_header(slide, "FEATURED PROJECT: DocQ", "03")
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.width = Pt(2)
    card.line.color.rgb = BLUE
    sum_txt = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11), Inches(1))
    p = sum_txt.text_frame.paragraphs[0]
    p.text = "DocQ: PDF-based Multiplayer 3D Quiz Game"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = NAVY
    p = sum_txt.text_frame.add_paragraph()
    p.text = "🏆 SSAFY 공통 프로젝트 반 우승 수상 | 2026.01.06 ~ 2026.02.13 (6주)"
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Save
    prs.save("portfolio-14.pptx")
    print("Success: portfolio-14.pptx created with chosen slogan and 3 core values.")

if __name__ == "__main__":
    create_final_slogan_ppt()
