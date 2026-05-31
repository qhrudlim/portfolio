from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy

def set_text_in_shape(shape, text):
    if hasattr(shape, "text_frame"):
        shape.text_frame.text = text

def duplicate_slide(prs, index):
    """Duplicates a slide by index."""
    template = prs.slides[index]
    blank_layout = prs.slide_layouts[6] # Blank layout
    new_slide = prs.slides.add_slide(blank_layout)
    
    for shape in template.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    
    return new_slide

def create_portfolio_7():
    prs = Presentation("portfolio-4.pptx")
    
    # Slide 1: Title
    slide1 = prs.slides[0]
    # Shapes: TextBox 3 (Frontend Developer...), TextBox 4 (임보경 포트폴리오...), TextBox 5 (User Experience...)
    for shape in slide1.shapes:
        if "Frontend Developer" in shape.text:
            shape.text_frame.text = "Frontend Developer"
        elif "임보경" in shape.text:
            shape.text_frame.text = "임보경 포트폴리오"
        elif "User Experience" in shape.text or "사용자 경험" in shape.text:
            shape.text_frame.text = "사용자 경험의 흐름을 논리적으로 설계하고,\n지속 가능한 코드 구조를 지향합니다."

    # Slide 2: ABOUT ME
    slide2 = prs.slides[1]
    # Update intro text if needed, but it seems already okay in portfolio-4.
    
    # Slide 3: TECH STACK
    slide3 = prs.slides[2]
    # It has 4 boxes: Languages, Frontend, 3D & Graphics, Tools.
    # The extracted info matches closely.

    # Slide 4: DocQ Overview
    slide4 = prs.slides[3]
    # Shape 5 (Description): PDF 기반 멀티플레이 3D 퀴즈 보드게임
    # Shape 6 (Details): Period, Personnel, Role
    for shape in slide4.shapes:
        if "PDF 기반" in shape.text:
            shape.text_frame.text = "DocQ – PDF 기반 멀티플레이 3D 퀴즈 보드게임"
        elif "기간:" in shape.text:
            shape.text_frame.text = (
                "기간: 2026.01.06 ~ 2026.02.13 (6주)\n\n"
                "인원: 6명 (프론트엔드 3명)\n\n"
                "역할: 프론트엔드 – 3D 게임 화면 설계 및 상태 기반 흐름 구조 설계"
            )

    # Slide 5: Technical Challenges (1)
    slide5 = prs.slides[4]
    # Challenge 1: 2.5D vs 3D
    for shape in slide5.shapes:
        if "Challenge: 2.5D" in shape.text:
            shape.text_frame.text = "Challenge 1: 2.5D 혼합 렌더링 한계"
        if "[Problem]" in shape.text and "렌더링 순서" in shape.text:
            shape.text_frame.text = (
                "[Problem]\n"
                "TresCanvas 내 2.5D(2D 맵 + 3D 캐릭터) 구성 시 렌더링 오류 발생.\n"
                "3D 맵 + 2D 캐릭터 구성 시 캐릭터 비출력 등 기술적 한계 확인.\n\n"
                "[Solution]\n"
                "3D 통합 구조로 전환 결정. 코드 기반 3D 맵으로 MVP를 우선 제작하여 "
                "웹 환경 구현 가능성을 조기 검증함."
            )

    # Slide 6: System Design & Result -> Rename to Technical Challenges (2)
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if "05. SYSTEM DESIGN" in shape.text:
            shape.text_frame.text = "05. TECHNICAL CHALLENGES (2)"
        # We'll replace the content with FBX vs GLTF info
        if "System:" in shape.text:
            shape.text_frame.text = (
                "Challenge 2: 에셋 파이프라인 및 애니메이션\n\n"
                "• FBX 포맷의 대용량으로 인한 초기 로딩 지연 문제 해결을 위해 GLTF/GLB로 전환.\n"
                "• Mixamo 애니메이션의 칸별 끊김 문제를 해결하기 위해 Blender에서 타임라인 분할 및 재구성."
            )
        if "Project Results" in shape.text:
            shape.text_frame.text = "Key Implementation"
        if "✔" in shape.text:
            shape.text_frame.text = (
                "✔ TresJS 기반 3D 환경 구성 및 인터랙션 구현\n"
                "✔ 룰렛 결과값과 캐릭터 이동 로직의 상태 동기화\n"
                "✔ 버튼 개입 없는 상태 전이 기반 자동 진행 구조"
            )

    # Add Slide 7: System Design & Result (Detailed)
    # We'll copy slide 6 as a template
    slide7 = duplicate_slide(prs, 5)
    for shape in slide7.shapes:
        if "05. TECHNICAL CHALLENGES" in shape.text:
            shape.text_frame.text = "06. SYSTEM DESIGN & RESULT"
        if "Challenge 2:" in shape.text:
            shape.text_frame.text = (
                "System: 상태 전이(State Transition) 기반 설계\n\n"
                "• 단계 간 의존성을 최소화하여 멀티플레이 환경에서의 동기화 로직을 안정화.\n"
                "• 이벤트 중심에서 상태 중심 설계로 전환하여 유지보수성 향상."
            )
        if "Key Implementation" in shape.text:
            shape.text_frame.text = "Project Results"
        if "✔" in shape.text:
            shape.text_frame.text = (
                "✔ 웹 환경에서의 3D 기반 게임 흐름 구현 가능성 검증\n"
                "✔ 에셋 파이프라인 최적화를 통한 팀 협업 효율 증대\n"
                "✔ 프로젝트 이후 React/Vanilla CSS 전환 구조 실험 완료"
            )

    # Add Slide 8: Other Projects
    slide8 = duplicate_slide(prs, 5) # Use slide 6 template again
    for shape in slide8.shapes:
        if "05. TECHNICAL CHALLENGES" in shape.text:
            shape.text_frame.text = "07. OTHER PROJECTS"
        if "Challenge 2:" in shape.text:
            shape.text_frame.text = (
                "BioTwin (SSAFY 특화 프로젝트)\n"
                "• 디지털 트윈 기술 기반 도메인 프로젝트 수행\n\n"
                "JI[1]PCHAK (SSAFY 자율 프로젝트)\n"
                "• 메인 페이지 및 GUIDE 구현, 하드웨어 제작 참여"
            )
        if "Key Implementation" in shape.text:
            shape.text_frame.text = "Experience"
        if "✔" in shape.text:
            shape.text_frame.text = (
                "✔ SSAFY 알고리즘 스터디 (Data Structure, DP 집중)\n"
                "✔ SWEA & BOJ 문제 해결 기록 다수 보유"
            )

    # Add Slide 9: Contact
    slide9 = duplicate_slide(prs, 0) # Use title slide as base for closing
    for shape in slide9.shapes:
        if "Frontend Developer" in shape.text:
            shape.text_frame.text = "CONTACT"
        elif "임보경" in shape.text:
            shape.text_frame.text = "감사합니다"
        elif "사용자 경험" in shape.text:
            shape.text_frame.text = (
                "Email: qhrudlim@gmail.com\n"
                "LinkedIn: linkedin.com/in/qhrudlim\n"
                "GitHub: github.com/qhrudlim"
            )

    prs.save("portfolio-7.pptx")
    print("Successfully created portfolio-7.pptx")

if __name__ == "__main__":
    create_portfolio_7()
