from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy

def set_text_in_shape(shape, text):
    if hasattr(shape, "text_frame"):
        shape.text_frame.text = text

def duplicate_slide(prs, index):
    """Duplicates a slide by index."""
    template = prs.slides[index]
    blank_layout = prs.slide_layouts[6] # Blank layout
    new_slide = prs.slides.add_slide(blank_layout)
    
    for shape in template.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    
    return new_slide

def create_portfolio_8():
    prs = Presentation("portfolio-4.pptx")
    
    # Slide 1: Intro (Modified from Template Slide 0)
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if "Frontend Developer" in shape.text:
            shape.text_frame.text = "Frontend Developer"
        elif "임보경" in shape.text:
            shape.text_frame.text = "임보경 포트폴리오"
        elif "User Experience" in shape.text or "사용자 경험" in shape.text:
            shape.text_frame.text = "사용자 경험의 흐름을 설계하고\n구조를 단순하게 유지하는 프론트엔드 개발자"

    # Slide 2: ABOUT ME & TECH STACK (Modified from Template Slide 1/2)
    slide2 = prs.slides[1]
    # Update About Me section
    for shape in slide2.shapes:
        if "ABOUT ME" in shape.text.upper():
            pass # Keep title
        elif "구조" in shape.text or "문제를" in shape.text:
            shape.text_frame.text = (
                "• 문제를 구조로 정리하는 것을 선호\n"
                "• 구현 이후 리팩토링과 UX 정제에 집중\n"
                "• 질문 전 충분한 고민과 자기 주도적 탐구 수행\n"
                "• 기술 전환을 통해 구조 이해도를 점검하는 습관 보유"
            )

    # Slide 3: DocQ Overview (Modified from Template Slide 3)
    slide3 = prs.slides[3]
    for shape in slide3.shapes:
        if "PDF 기반" in shape.text or "DocQ" in shape.text:
            shape.text_frame.text = "DocQ – PDF 기반 멀티플레이 3D 퀴즈 게임"
        elif "기간:" in shape.text:
            shape.text_frame.text = (
                "기간: 2026.01.06 ~ 2026.02.13 (6주)\n"
                "인원: 6명 (프론트엔드 3명)\n"
                "역할: 프론트엔드 – 3D 게임 화면 및 상태 기반 흐름 설계\n\n"
                "개요: PDF 기반 퀴즈 생성 및 3D 보드게임 형태의 멀티플레이 서비스"
            )

    # Slide 4: Technical Challenges (1) - 2.5D vs 3D
    slide4 = prs.slides[4]
    for shape in slide4.shapes:
        if "Challenge:" in shape.text:
            shape.text_frame.text = "Challenge 1: 2.5D 혼합 렌더링 한계"
        if "[Problem]" in shape.text:
            shape.text_frame.text = (
                "[Problem]\n"
                "• TresCanvas 내 2.5D(2D 맵 + 3D 캐릭터) 구성 시 렌더링 오류 발생\n"
                "• 3D 맵 + 2D 캐릭터 구성 시 캐릭터 비출력 문제 확인\n\n"
                "[Solution]\n"
                "• 3D 통합 구조로 전환 결정\n"
                "• 코드 기반 3D 맵으로 MVP 제작하여 구현 가능성 조기 검증"
            )

    # Slide 5: Technical Challenges (2) - Asset Pipeline
    slide5 = prs.slides[5]
    for shape in slide5.shapes:
        if "05. SYSTEM DESIGN" in shape.text:
            shape.text_frame.text = "TECHNICAL CHALLENGES (2)"
        if "System:" in shape.text:
            shape.text_frame.text = (
                "Challenge 2: 에셋 파이프라인 최적화\n\n"
                "• FBX 포맷의 대용량으로 인한 초기 로딩 지연 -> GLTF/GLB 전환\n"
                "• Mixamo 애니메이션의 칸별 끊김 현상 -> Blender에서 타임라인 재구성\n"
                "• FBX → GLTF 변환 시 메시 오류 해결 (export 옵션 조정)"
            )
        if "Project Results" in shape.text:
            shape.text_frame.text = "Key Implementation"
        if "✔" in shape.text:
            shape.text_frame.text = (
                "✔ TresJS 기반 3D 환경 구성 및 인터랙션 구현\n"
                "✔ 캐릭터 애니메이션 타임라인 최적화\n"
                "✔ 웹 환경에 최적화된 에셋 파이프라인 구축"
            )

    # Slide 6: System Design & Logic (Duplicate from Slide 5)
    slide6 = duplicate_slide(prs, 5)
    for shape in slide6.shapes:
        if "05. SYSTEM DESIGN" in shape.text or "TECHNICAL" in shape.text:
            shape.text_frame.text = "SYSTEM DESIGN & LOGIC"
        if "Challenge 2:" in shape.text:
            shape.text_frame.text = (
                "상태 전이(State Transition) 기반 설계\n\n"
                "• 문제 -> 정답 -> 룰렛 -> 이동 -> 미니게임의 자동 진행 구조\n"
                "• 이벤트 중심에서 상태 중심 설계로 전환하여 단계 간 의존성 최소화\n"
                "• 멀티플레이 환경에서의 단계 동기화 안정성 확보"
            )
        if "Key Implementation" in shape.text:
            shape.text_frame.text = "Flow"
        if "✔" in shape.text:
            shape.text_frame.text = (
                "✔ 상태 전이 기반 자동 진행 로직\n"
                "✔ 룰렛 결과값과 캐릭터 이동 로직의 동기화\n"
                "✔ 컴포넌트 책임 분리 및 코드 리팩토링"
            )

    # Slide 7: Result & Retrospective (Duplicate from Slide 5)
    slide7 = duplicate_slide(prs, 5)
    for shape in slide7.shapes:
        if "05. SYSTEM DESIGN" in shape.text or "TECHNICAL" in shape.text:
            shape.text_frame.text = "RESULT & RETROSPECTIVE"
        if "Challenge 2:" in shape.text:
            shape.text_frame.text = (
                "성과 및 배운 점\n\n"
                "• 웹 환경에서 3D 기반 멀티플레이 게임 구현 가능성 검증\n"
                "• 실험 기반 의사결정(2.5D -> 3D)의 중요성 체감\n"
                "• 프로젝트 이후 React/Vanilla CSS 전환 구조 실험 완료"
            )
        if "Key Implementation" in shape.text:
            shape.text_frame.text = "Conclusion"
        if "✔" in shape.text:
            shape.text_frame.text = (
                "✔ 3D 에셋 포맷 선택 및 변환 과정의 중요성 확인\n"
                "✔ 초기 구조 설계가 프로젝트 전체에 미치는 영향 학습\n"
                "✔ 사용자 경험(UX) 디테일 개선의 가치 발견"
            )

    # Slide 8: Other Projects & Contact
    slide8 = duplicate_slide(prs, 0)
    for shape in slide8.shapes:
        if "Frontend Developer" in shape.text:
            shape.text_frame.text = "CONTACT & OTHERS"
        elif "임보경" in shape.text:
            shape.text_frame.text = "감사합니다"
        elif "사용자 경험" in shape.text:
            shape.text_frame.text = (
                "Others: BioTwin (디지털 트윈), JIPCHAK (자율 프로젝트)\n\n"
                "Email: qhrudlim@gmail.com\n"
                "LinkedIn: linkedin.com/in/qhrudlim\n"
                "GitHub: github.com/qhrudlim"
            )

    prs.save("portfolio-8.pptx")
    print("Successfully created portfolio-8.pptx")

if __name__ == "__main__":
    create_portfolio_8()
