from pptx import Presentation

def inspect_pptx(file_path):
    prs = Presentation(file_path)
    print(f"File: {file_path}")
    print(f"Number of slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i+1}:")
        for shape in slide.shapes:
            print(f"  Shape Type: {shape.shape_type}, Name: {shape.name}")
            if hasattr(shape, "text"):
                print(f"    Text: {shape.text[:50]}...")
            if shape.has_table:
                print("    (Table)")
            if shape.has_chart:
                print("    (Chart)")

inspect_pptx("portfolio-4.pptx")
