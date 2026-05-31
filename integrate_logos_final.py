from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def integrate_real_logos():
    # Load version 14 as base
    try:
        prs = Presentation('portfolio-14.pptx')
    except Exception as e:
        print(f"Error loading portfolio-14.pptx: {e}")
        return

    # Slide 3 (Index 2)
    slide = prs.slides[2]

    # professional palette
    BLUE = RGBColor(52, 152, 219)
    
    # Positioning for Logos at the bottom of the featured project
    # Adjusted to fit next to or under the challenges box
    x_start = Inches(0.8)
    y_start = Inches(6.0)
    logo_size = Inches(0.8)
    
    # Label for logos
    label_box = slide.shapes.add_textbox(x_start, y_start - Inches(0.35), Inches(3), Inches(0.3))
    p = label_box.text_frame.paragraphs[0]
    p.text = "CORE TECHNOLOGIES"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = BLUE

    # Technology mapping
    tech_logos = [
        ("logos/react.png", "React"),
        ("logos/threejs.png", "Three.js"),
        ("logos/typescript.png", "TypeScript"),
        ("logos/blender.png", "Blender")
    ]

    for i, (path, name) in enumerate(tech_logos):
        if os.path.exists(path):
            curr_x = x_start + Inches(i * 1.5)
            # Add Picture
            slide.shapes.add_picture(path, curr_x, y_start, width=logo_size)
            # Optional: Add small text under logo
            txt_box = slide.shapes.add_textbox(curr_x, y_start + logo_size, logo_size, Inches(0.3))
            p = txt_box.text_frame.paragraphs[0]
            p.text = name
            p.font.size = Pt(10)
            p.font.bold = True
            p.alignment = 1 # Center

    prs.save("portfolio-15.pptx")
    print("Success: portfolio-15.pptx created with actual tech logo images.")

if __name__ == "__main__":
    integrate_real_logos()
