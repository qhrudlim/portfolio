from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_logo_integrated_ppt():
    # Load version 14 as base
    try:
        prs = Presentation('portfolio-14.pptx')
    except:
        print("Error: portfolio-14.pptx not found. Creating from script logic.")
        return

    # Assuming Slide 3 (Index 2) is the Projects/Milestones page
    # Based on previous analysis: Slide 1 (Cover), Slide 2 (About/Skills), Slide 3 (Projects/Milestones)
    slide = prs.slides[2]

    # professional palette
    NAVY = RGBColor(26, 37, 47)      
    BLUE = RGBColor(52, 152, 219)    
    TEXT_GRAY = RGBColor(80, 80, 80)
    GOLD = RGBColor(241, 196, 15)

    # We will look for tech stack text on Slide 3 and overlay/replace with a small icon section
    # Since I cannot 'see' the manual changes made by user in pptx, 
    # I will add a fresh, styled tech logo row to Slide 3 based on DocQ and other projects.
    
    # Adding a visual tech stack area for DocQ on Slide 3
    # Coordinates based on Slide 3 layout (Featured Project area)
    x_start = Inches(0.8)
    y_start = Inches(4.5) # Assuming space under challenges or next to it
    
    # Text Label
    label = slide.shapes.add_textbox(x_start, y_start - Inches(0.4), Inches(3), Inches(0.4))
    p = label.text_frame.paragraphs[0]
    p.text = "CORE STACK"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = BLUE

    # In a real scenario, I'd insert image files. 
    # Since I don't have individual tech logo images as files, I will use high-quality styled shapes 
    # and labels to represent the 'Logo' look requested.
    techs = ["React", "Three.js", "TS", "Blender"]
    
    for i, tech in enumerate(techs):
        tx = x_start + Inches(i * 1.5)
        # Decorative Icon Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tx, y_start, Inches(1.2), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE if i % 2 == 0 else BLUE_LIGHT # Placeholder logic for contrast
        box.line.color.rgb = BLUE
        box.line.width = Pt(1)
        
        # Tech Name (as logo replacement)
        txt = slide.shapes.add_textbox(tx, y_start, Inches(1.2), Inches(0.6))
        tp = txt.text_frame.paragraphs[0]
        tp.text = tech
        tp.alignment = PP_ALIGN.CENTER
        tp.font.bold = True
        tp.font.size = Pt(12)
        tp.font.color.rgb = NAVY

    prs.save("portfolio-15.pptx")
    print("Success: portfolio-15.pptx created with integrated tech stack visual blocks.")

if __name__ == "__main__":
    # Standardizing logic: 
    # Since the user manually edited 14, I should use the python-pptx library to 
    # modify the existing slide objects instead of re-creating from scratch.
    
    prs = Presentation('portfolio-14.pptx')
    slide = prs.slides[2] # Slide 3
    
    # Colors
    NAVY = RGBColor(26, 37, 47)      
    BLUE = RGBColor(52, 152, 219)
    WHITE = RGBColor(255, 255, 255)
    
    # Search for text 'React' or '기술' to replace
    # But more safely, add a dedicated visual row for the tech logos
    # Positioning this in the project details section
    
    y_pos = Inches(5.8) # Lower part of Slide 3
    x_pos = Inches(1.0)
    
    # Adding a header for tech
    header = slide.shapes.add_textbox(x_pos, y_pos - Inches(0.3), Inches(3), Inches(0.4))
    hp = header.text_frame.paragraphs[0]
    hp.text = "TECHNOLOGIES"
    hp.font.bold = True
    hp.font.size = Pt(12)
    hp.font.color.rgb = BLUE
    
    # Adding logo-like blocks
    techs = [("React", "#61DAFB"), ("Three.js", "#000000"), ("TS", "#3178C6"), ("Blender", "#F5792A")]
    
    for i, (name, color_hex) in enumerate(techs):
        tx = x_pos + Inches(i * 1.3)
        # The 'Logo' Box
        logo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tx, y_pos, Inches(1.1), Inches(0.8))
        logo_box.fill.solid()
        logo_box.fill.fore_color.rgb = WHITE
        logo_box.line.color.rgb = BLUE
        
        # Name inside
        nt = slide.shapes.add_textbox(tx, y_pos + Inches(0.2), Inches(1.1), Inches(0.4))
        np = nt.text_frame.paragraphs[0]
        np.text = name
        np.alignment = PP_ALIGN.CENTER
        np.font.bold = True
        np.font.size = Pt(11)
        np.font.color.rgb = NAVY

    prs.save("portfolio-15.pptx")
    print("Success: portfolio-15.pptx updated with visual tech blocks on Slide 3.")
